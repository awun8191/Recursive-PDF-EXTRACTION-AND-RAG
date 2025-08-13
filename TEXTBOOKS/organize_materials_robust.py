#!/usr/bin/env python3
"""
Hierarchical Course Document Classifier
--------------------------------------

Classifies documents from an input folder into a **copy-only** mirror under an
output root using this hierarchy:

    <output-root>/<DEPT_CODE>/<LEVEL>/<SEM_NUM>/<COURSE_CODE>/<filename>

It reads `courses.json` (your engineering college outline) and supports fast
TF‑IDF similarity or optional Sentence‑Transformers (SBERT). It extracts text
from the first N pages of PDF/DOCX/PPTX/TXT and optionally OCRs scanned PDFs.

Robust JSON ingestion: accepts a list of course dicts, a mapping of codes to
course dicts, or nested structures like `{DEPT->{LEVEL->{SEM->[courses]}}}`.

Usage example
-------------
python classifier.py \
  --input-dir /path/incoming \
  --courses-json /path/courses.json \
  --output-root /path/SortedDocs \
  --pages 4 \
  --method tfidf \
  --threshold 0.30 \
  --workers 6 \
  --layout hierarchical

Install
-------
  pip install pypdf python-docx python-pptx scikit-learn
  # optional
  pip install sentence-transformers pillow pytesseract pdf2image google-generativeai

Notes
-----
- **Copy by default** (originals untouched). Use `--move` to move instead.
- Uncertain files go to `<output-root>/_unclassified/`.
- Semester strings FIRST/SECOND map to 1/2.
- Caches course vectors under `<output-root>/.cache/`.

License: MIT
"""
from __future__ import annotations

import argparse
import concurrent.futures
import dataclasses
import functools
import hashlib
import json
import logging
import os
import re
import shutil
import sys
import time
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

# --- Lazy optional imports ---
try:
    from pypdf import PdfReader  # type: ignore
    PYPDF_AVAILABLE = True
except Exception:
    PYPDF_AVAILABLE = False

try:
    import docx  # python-docx  # type: ignore
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation  # python-pptx  # type: ignore
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    import pytesseract  # type: ignore
    PYTESSERACT_AVAILABLE = True
except Exception:
    PYTESSERACT_AVAILABLE = False

try:
    from pdf2image import convert_from_path  # type: ignore
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False

try:
    import google.generativeai as genai  # type: ignore
    GEMINI_AVAILABLE = True
except Exception:
    GEMINI_AVAILABLE = False

from sklearn.feature_extraction.text import TfidfVectorizer  # type: ignore
from sklearn.metrics.pairwise import cosine_similarity  # type: ignore

SBERT_AVAILABLE = False
try:
    from sentence_transformers import SentenceTransformer  # type: ignore
    import numpy as np  # type: ignore
    SBERT_AVAILABLE = True
except Exception:
    pass

# --- Logging ---
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger("doc_classifier")

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt"}

# -------------------------
# JSON helpers & normalizer
# -------------------------

def read_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def write_json(path: Path, obj: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        json.dump(obj, f, ensure_ascii=False, indent=2)


def normalize_courses(courses_raw: Any) -> List[Dict[str, Any]]:
    """Accept many shapes and return a flat list of course dicts with at least:
    code, title, department_code, level, semester
    """
    out: List[Dict[str, Any]] = []

    def safe(val: Any) -> Optional[str]:
        return None if val is None else str(val)

    def coerce_course(x: Any, hints: Dict[str, Any] = None):
        hints = hints or {}
        if isinstance(x, dict):
            c = dict(x)  # shallow copy
            # Fill missing fields from hints
            for k in ("department_code", "level", "semester"):
                c.setdefault(k, hints.get(k))
            # Normalize
            if c.get("code"):
                c["code"] = str(c["code"]).strip().upper()
            if c.get("title"):
                c["title"] = str(c["title"]).strip()
            if c.get("department_code"):
                c["department_code"] = str(c["department_code"]).strip().upper()
            if c.get("level"):
                # keep numeric level like "300"
                lev = re.sub(r"\D+", "", str(c["level"]))
                c["level"] = lev or str(c["level"]).strip()
            if c.get("semester"):
                c["semester"] = str(c["semester"]).strip()
            out.append(c)
        elif isinstance(x, str):
            s = x.strip()
            m = re.match(r"^([A-Za-z]{2,})\s*-?\s*(\d{2,3})\b(.*)$", s)
            code = title = None
            dept = level = sem = None
            if m:
                dept = m.group(1).upper()
                level = m.group(2)[0] + "00" if len(m.group(2)) >= 3 else None
                code = f"{dept} {m.group(2)}"
                tail = m.group(3).strip(" -")
                title = tail
                if "FIRST" in s.upper():
                    sem = "FIRST"
                elif "SECOND" in s.upper():
                    sem = "SECOND"
            c = {
                "code": code or s.upper(),
                "title": title or "",
                "department_code": hints.get("department_code") or dept,
                "level": hints.get("level") or level,
                "semester": hints.get("semester") or sem,
            }
            out.append(c)
        else:
            # Unknown type; ignore
            pass

    def walk(node: Any, hints: Dict[str, Any] = None):
        hints = hints or {}
        if isinstance(node, list):
            for itm in node:
                walk(itm, hints)
        elif isinstance(node, dict):
            # Direct course dict
            if all(k in node for k in ("code", "title")):
                coerce_course(node, hints)
                return
            # String-keyed dict: could be nested structure
            if all(isinstance(k, str) for k in node.keys()):
                for k, v in node.items():
                    keyu = str(k).upper().strip()
                    if keyu in {"FIRST", "SECOND", "1", "2"}:
                        walk(v, {**hints, "semester": "FIRST" if keyu in {"FIRST","1"} else "SECOND"})
                    elif re.fullmatch(r"\d{3}", keyu):
                        walk(v, {**hints, "level": keyu})
                    elif re.fullmatch(r"[A-Z]{2,5}", keyu):
                        walk(v, {**hints, "department_code": keyu})
                    else:
                        if keyu in {"COURSES", "LIST", "OUTLINE", "CATALOG"}:
                            walk(v, hints)
                        else:
                            if isinstance(v, dict) and ("code" in v or "title" in v):
                                vv = dict(v)
                                vv.setdefault("code", k)
                                coerce_course(vv, hints)
                            else:
                                walk(v, hints)
            else:
                for v in node.values():
                    walk(v, hints)
        else:
            coerce_course(node, hints)

    walk(courses_raw)

    # Cleanup and best-effort derivations
    cleaned = []
    for c in out:
        if not c.get("code"):
            continue
        # Derive dept from code if missing
        if not c.get("department_code"):
            m = re.match(r"^([A-Za-z]{2,})\s*-?\s*(\d{2,3})", c["code"])
            if m:
                c["department_code"] = m.group(1).upper()
        cleaned.append(c)
    return cleaned


# --------------------
# Utility & formatting
# --------------------

def file_sha256(path: Path, chunk: int = 1 << 20) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        while True:
            b = f.read(chunk)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


def normalize_space(s: str) -> str:
    s = s.replace("\x00", " ")
    s = re.sub(r"\s+", " ", s)
    return s.strip()


def safe_segment(name: str) -> str:
    name = re.sub(r"[\\/]+", " ", str(name)).strip()
    name = re.sub(r"[^A-Za-z0-9 _.-]+", "", name)
    name = re.sub(r"\s+", " ", name)
    return name


def semester_to_number(sem: Optional[str]) -> str:
    if not sem:
        return "1"
    s = str(sem).strip().upper()
    if s in {"1", "FIRST", "FIRST SEMESTER", "SEMESTER 1"}:
        return "1"
    if s in {"2", "SECOND", "SECOND SEMESTER", "SEMESTER 2"}:
        return "2"
    return "1"


# -----------------
# Course corpus API
# -----------------

def course_to_text(course: Dict[str, Any]) -> str:
    parts: List[str] = []
    for key in ("code", "title", "department", "department_code", "level", "semester"):
        val = course.get(key)
        if val:
            parts.append(str(val))
    desc = course.get("description")
    if desc:
        parts.append(desc)
    topics = course.get("topics") or []
    if isinstance(topics, list):
        parts.extend([str(t) for t in topics])
    elif isinstance(topics, str):
        parts.append(topics)
    return normalize_space(" ".join(parts))


def build_course_corpus(courses: List[Dict[str, Any]]) -> Tuple[List[str], List[str]]:
    labels: List[str] = []
    texts: List[str] = []
    for c in courses:
        if not isinstance(c, dict):
            continue
        code = str(c.get("code") or "").strip()
        if not code:
            continue
        labels.append(code)
        texts.append(course_to_text(c))
    return labels, texts


def build_code_meta_map(courses: List[Dict[str, Any]]) -> Dict[str, Dict[str, str]]:
    m: Dict[str, Dict[str, str]] = {}
    for c in courses:
        code = str(c.get("code") or "").strip()
        if not code:
            continue
        dept_code = c.get("department_code") or c.get("dept_code") or ""
        level = c.get("level") or ""
        sem = c.get("semester") or ""
        m[code] = {
            "department_code": safe_segment(dept_code) or "UNKNOWN",
            "level": safe_segment(str(level)) or "000",
            "semester_num": semester_to_number(sem),
            "course_code": safe_segment(code),
        }
    return m


# ----------------
# Text extraction
# ----------------

def extract_text_pdf(path: Path, max_pages: int, ocr: bool = False, ocr_dpi: int = 250) -> str:
    text_parts: List[str] = []
    if PYPDF_AVAILABLE:
        try:
            reader = PdfReader(str(path))
            pages = min(len(reader.pages), max_pages)
            for i in range(pages):
                try:
                    txt = reader.pages[i].extract_text() or ""
                    if txt.strip():
                        text_parts.append(txt)
                except Exception as e:
                    logger.debug("PDF text extraction failed on page %s: %s", i, e)
        except Exception as e:
            logger.debug("PDF open failed via pypdf for %s: %s", path, e)

    if ocr and (not text_parts or sum(len(t) for t in text_parts) < 40):
        if not (PDF2IMAGE_AVAILABLE and PYTESSERACT_AVAILABLE):
            logger.warning("OCR requested but pdf2image/pytesseract not available for %s", path.name)
        else:
            try:
                images = convert_from_path(str(path), dpi=ocr_dpi, first_page=1, last_page=max_pages)
                for img in images:
                    try:
                        txt = pytesseract.image_to_string(img)
                        if txt.strip():
                            text_parts.append(txt)
                    except Exception as e:
                        logger.debug("Tesseract failed: %s", e)
            except Exception as e:
                logger.debug("pdf2image failed: %s", e)

    return normalize_space(" ".join(text_parts))


def extract_text_docx(path: Path, max_pages: int) -> str:
    if not DOCX_AVAILABLE:
        return ""
    try:
        doc = docx.Document(str(path))
        para_texts = []
        chars = 0
        for p in doc.paragraphs:
            t = p.text.strip()
            if not t:
                continue
            para_texts.append(t)
            chars += len(t)
            if chars > max_pages * 1800:
                break
        return normalize_space(" ".join(para_texts))
    except Exception:
        return ""


def extract_text_pptx(path: Path, max_pages: int) -> str:
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(str(path))
        texts = []
        for i, slide in enumerate(prs.slides):
            if i >= max_pages:
                break
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    t = shp.text.strip()
                    if t:
                        texts.append(t)
        return normalize_space(" ".join(texts))
    except Exception:
        return ""


def extract_text_txt(path: Path, max_chars: int = 12000) -> str:
    try:
        with path.open("r", encoding="utf-8", errors="ignore") as f:
            return normalize_space(f.read(max_chars))
    except Exception:
        return ""


def extract_text(path: Path, pages: int, ocr: bool) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_text_pdf(path, pages, ocr=ocr)
    elif ext == ".docx":
        return extract_text_docx(path, pages)
    elif ext == ".pptx":
        return extract_text_pptx(path, pages)
    elif ext == ".txt":
        return extract_text_txt(path)
    else:
        return ""


# --------------------------
# Vectorization & LLM refine
# --------------------------

@dataclasses.dataclass
class VectorIndex:
    method: str  # "tfidf" or "sbert"
    labels: List[str]
    model_name: Optional[str] = None
    tfidf: Optional[TfidfVectorizer] = None
    tfidf_matrix: Optional[Any] = None
    sbert_model: Optional[Any] = None
    sbert_vectors: Optional[Any] = None


def build_index_tfidf(course_texts: List[str]) -> Tuple[TfidfVectorizer, Any]:
    vec = TfidfVectorizer(
        lowercase=True,
        stop_words="english",
        ngram_range=(1, 2),
        min_df=1,
        max_df=0.95
    )
    X = vec.fit_transform(course_texts)
    return vec, X


def build_index_sbert(course_texts: List[str], model_name: str = "all-MiniLM-L6-v2"):
    if not SBERT_AVAILABLE:
        raise RuntimeError("sentence-transformers is not installed; cannot use --method sbert")
    model = SentenceTransformer(model_name)
    vectors = model.encode(course_texts, batch_size=32, show_progress_bar=False, normalize_embeddings=True)
    return model, vectors


def llm_refine_gemini(candidates: List[Tuple[str, float]], snippet: str, course_map: Dict[str, str]) -> Optional[str]:
    api_key = "AIzaSyCP6igfyX0FTLiWxN0os50nvN748gn6YiA"
    if not (GEMINI_AVAILABLE and api_key):
        return None
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel("models/gemini-1.5-flash")
        top_list = "\n".join([f"{c[0]}: {course_map.get(c[0], '')[:200]}" for c in candidates])
        prompt = f"""You are a course document classifier.
Pick the SINGLE best matching course code from the shortlist strictly, or respond "NONE" if not sure.

SNIPPET (truncated):
{snippet[:2500]}

SHORTLIST:
{top_list}

Reply with only the course code or NONE."""
        resp = model.generate_content(prompt)
        choice = (resp.text or "").strip()
        if choice.upper() == "NONE":
            return None
        valid_codes = {c[0] for c in candidates}
        if choice in valid_codes:
            return choice
        for c in valid_codes:
            if choice.upper() == c.upper():
                return c
    except Exception:
        return None
    return None


# ---------------
# Classification
# ---------------

@dataclasses.dataclass
class ClassificationResult:
    file: Path
    assigned: Optional[str]
    confidence: float
    topk: List[Tuple[str, float]]
    reason: str


def classify_one(
    path: Path,
    pages: int,
    ocr: bool,
    index: VectorIndex,
    course_texts: List[str],
    label_to_text: Dict[str, str],
    threshold: float,
    use_llm: bool,
    top_k: int = 5
) -> ClassificationResult:
    text = extract_text(path, pages, ocr=ocr)
    if not text:
        return ClassificationResult(path, None, 0.0, [], "no_text")

    if index.method == "tfidf":
        assert index.tfidf is not None and index.tfidf_matrix is not None
        q = index.tfidf.transform([text])
        sims = cosine_similarity(q, index.tfidf_matrix)[0]
    else:
        assert index.sbert_model is not None and index.sbert_vectors is not None
        qvec = index.sbert_model.encode([text], normalize_embeddings=True)[0]
        sims = (index.sbert_vectors @ qvec)

    ranked = sorted(zip(index.labels, sims), key=lambda x: x[1], reverse=True)
    topk = ranked[:top_k]
    best_code, best_score = topk[0]

    assigned: Optional[str] = None
    reason = "vector_match"
    if best_score >= threshold:
        assigned = best_code
    else:
        reason = f"below_threshold({best_score:.3f} < {threshold:.3f})"

    if use_llm and not assigned:
        llm_pick = llm_refine_gemini(topk, text, label_to_text)
        if llm_pick:
            assigned = llm_pick
            reason = "llm_refine"

    return ClassificationResult(path, assigned, float(best_score), topk, reason)


# ---------
# File ops
# ---------

def move_or_copy(src: Path, dst_dir: Path, move: bool) -> None:
    dst_dir.mkdir(parents=True, exist_ok=True)
    dst = dst_dir / src.name
    if dst.exists():
        prefix = file_sha256(src)[:8]
        dst = dst_dir / f"{prefix}_{src.name}"
    if move:
        shutil.move(str(src), str(dst))
    else:
        shutil.copy2(str(src), str(dst))


# --------------
# Cache handling
# --------------

def cache_paths(output_root: Path, method: str) -> Tuple[Path, Path]:
    meta = output_root / ".cache" / f"courses_{method}.meta.json"
    vecs = output_root / ".cache" / f"courses_{method}.npz"
    return meta, vecs


def load_index_from_cache(output_root: Path, method: str, model_name: Optional[str], course_texts: List[str], course_labels: List[str]) -> Optional[VectorIndex]:
    meta_path, vecs_path = cache_paths(output_root, method)
    if not meta_path.exists():
        return None
    try:
        with meta_path.open("r", encoding="utf-8") as f:
            meta = json.load(f)
        if meta.get("method") != method:
            return None
        if model_name and meta.get("model_name") != model_name:
            return None
        if meta.get("labels") != course_labels:
            return None
        index = VectorIndex(method=method, labels=course_labels, model_name=model_name)
        if method == "tfidf":
            tfidf, X = build_index_tfidf(course_texts)
            index.tfidf = tfidf
            index.tfidf_matrix = X
            return index
        else:
            if not SBERT_AVAILABLE:
                return None
            import numpy as _np  # type: ignore
            model, _ = build_index_sbert(course_texts, model_name=model_name or "all-MiniLM-L6-v2")
            arr = _np.load(vecs_path)["vectors"]
            index.sbert_model = model
            index.sbert_vectors = arr
            return index
    except Exception:
        return None


def save_index_to_cache(output_root: Path, index: VectorIndex, course_texts: List[str]) -> None:
    meta_path, vecs_path = cache_paths(output_root, index.method)
    meta = {"method": index.method, "labels": index.labels, "model_name": index.model_name}
    meta_path.parent.mkdir(parents=True, exist_ok=True)
    with meta_path.open("w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    if index.method == "sbert" and SBERT_AVAILABLE:
        import numpy as _np  # type: ignore
        if index.sbert_vectors is None:
            assert index.sbert_model is not None
            vecs = index.sbert_model.encode(course_texts, batch_size=32, show_progress_bar=False, normalize_embeddings=True)
        else:
            vecs = index.sbert_vectors
        _np.savez_compressed(vecs_path, vectors=vecs)


# -----
# CLI
# -----

def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Classify documents into course-code folders using course outlines.")
    p.add_argument("--input-dir", required=True, help="Folder containing incoming documents")
    p.add_argument("--courses-json", required=True, help="Path to courses.json")
    p.add_argument("--output-root", required=True, help="Root folder to mirror classified copies (originals untouched)")
    p.add_argument("--pages", type=int, default=4, help="Number of first pages to read")
    p.add_argument("--ocr", action="store_true", help="Enable OCR for scanned PDFs (requires pdf2image + pytesseract)")
    p.add_argument("--method", choices=["tfidf", "sbert"], default="tfidf", help="Vectorization method")
    p.add_argument("--sbert-model", default="all-MiniLM-L6-v2", help="Sentence-Transformers model name (if --method sbert)")
    p.add_argument("--threshold", type=float, default=0.30, help="Similarity threshold for auto-assign")
    p.add_argument("--llm", choices=["none", "gemini"], default="none", help="Optional LLM refinement (Gemini uses $GEMINI_API_KEY)")
    p.add_argument("--top-k", type=int, default=5, help="Keep top-K candidates for logs/LLM")
    p.add_argument("--move", action="store_true", help="Move files instead of copy (default is COPY)")
    p.add_argument("--workers", type=int, default=min(8, (os.cpu_count() or 4)), help="Parallel workers")
    p.add_argument("--dry-run", action="store_true", help="Do not move/copy files; just log decisions")
    p.add_argument("--include-exts", default=",".join(sorted(SUPPORTED_EXTS)), help="Comma-separated list of extensions to include")
    p.add_argument("--layout", choices=["hierarchical","flat"], default="hierarchical",
                   help="Folder layout: hierarchical = DEPT_CODE/LEVEL/SEM_NUM/COURSE_CODE (default), flat = COURSE_CODE only")
    return p.parse_args()


# -----
# Main
# -----

def main():
    args = parse_args()
    input_dir = Path(args.input_dir).expanduser().resolve()
    courses_json = Path(args.courses_json).expanduser().resolve()
    output_root = Path(args.output_root).expanduser().resolve()
    pages = max(1, args.pages)
    ocr = bool(args.ocr)
    method = args.method
    sbert_model = args.sbert_model
    threshold = args.threshold
    use_llm = (args.llm == "gemini")
    top_k = args.top_k
    move = bool(args.move)
    workers = max(1, int(args.workers))
    layout = args.layout

    include_exts = {e.strip().lower() if e.strip().startswith(".") else "." + e.strip().lower()
                    for e in args.include_exts.split(",") if e.strip()}

    # Read & normalize courses
    if not courses_json.exists():
        logger.error("courses.json not found at %s", courses_json)
        sys.exit(2)
    courses_raw = read_json(courses_json)
    courses: List[Dict[str, Any]] = normalize_courses(courses_raw)

    if not courses:
        logger.error("No courses could be parsed from %s", courses_json)
        sys.exit(2)

    logger.info("Loaded %d normalized courses (examples): %s",
                len(courses), [c.get("code") for c in courses[:5]])

    labels, texts = build_course_corpus(courses)
    label_to_text = {lbl: txt for lbl, txt in zip(labels, texts)}
    code_meta = build_code_meta_map(courses)

    # Build/load vector index
    index = load_index_from_cache(output_root, method, sbert_model if method == "sbert" else None, texts, labels)
    if index is None:
        logger.info("Building vector index (%s)...", method)
        index = VectorIndex(method=method, labels=labels, model_name=(sbert_model if method == "sbert" else None))
        if method == "tfidf":
            tfidf, X = build_index_tfidf(texts)
            index.tfidf = tfidf
            index.tfidf_matrix = X
        else:
            model, vecs = build_index_sbert(texts, model_name=sbert_model)
            index.sbert_model = model
            index.sbert_vectors = vecs
        save_index_to_cache(output_root, index, texts)
    else:
        logger.info("Loaded vector index from cache.")

    # Gather files
    files: List[Path] = []
    for root, _, fnames in os.walk(input_dir):
        for fn in fnames:
            p = Path(root) / fn
            if p.suffix.lower() in include_exts:
                files.append(p)
    files.sort()
    if not files:
        logger.warning("No input files found under %s (exts=%s)", input_dir, sorted(include_exts))
        return

    # Copy by default (originals untouched)
    if not move:
        logger.info("Copy mode: originals remain untouched.")

    unclassified_dir = output_root / "_unclassified"
    results: List[ClassificationResult] = []

    logger.info("Classifying %d file(s) with %d worker(s)...", len(files), workers)

    func = functools.partial(
        classify_one,
        pages=pages,
        ocr=ocr,
        index=index,
        course_texts=texts,
        label_to_text=label_to_text,
        threshold=threshold,
        use_llm=use_llm,
        top_k=top_k
    )

    start = time.time()
    with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as ex:
        for res in ex.map(func, files):
            results.append(res)
            assigned = res.assigned or "_unclassified"
            logger.info("FILE: %s  ->  %s  (%.3f)  reason=%s  top1..3=%s",
                        res.file.name, assigned, res.confidence, res.reason, res.topk[:3])

            if not args.dry_run:
                if res.assigned:
                    if layout == "hierarchical":
                        meta = code_meta.get(res.assigned, {})
                        dept = meta.get("department_code", "UNKNOWN") or "UNKNOWN"
                        level = meta.get("level", "000") or "000"
                        sem = meta.get("semester_num", "1") or "1"
                        course_code = meta.get("course_code", safe_segment(res.assigned))
                        target_dir = output_root / dept / level / sem / course_code
                    else:
                        target_dir = output_root / safe_segment(res.assigned)
                else:
                    target_dir = unclassified_dir

                try:
                    move_or_copy(res.file, target_dir, move=move)
                except Exception as e:
                    logger.error("Failed to move/copy %s: %s", res.file, e)

    elapsed = time.time() - start
    logger.info("Done in %.2fs. %d classified, %d unclassified.",
                elapsed,
                sum(1 for r in results if r.assigned),
                sum(1 for r in results if not r.assigned))

    # Report
    report = []
    for r in results:
        report.append({
            "file": str(r.file),
            "assigned": r.assigned,
            "confidence": r.confidence,
            "topk": [[c, float(s)] for c, s in r.topk],
            "reason": r.reason
        })
    write_json(output_root / "classification_report.json", report)


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        logger.warning("Interrupted by user")
        sys.exit(130)
