#!/usr/bin/env python3
"""
Classifier v5 — strict layout + shared 100/200 roots (no extra folders)
======================================================================

What this version guarantees (based on your constraints):

• **Only** these directory shapes are created:
    1)  <OUTPUT>/<DEPT>/<LEVEL>/<SEM>/<COURSE_CODE>/
    2)  <OUTPUT>/<SHARED_LEVEL>/<SEM>/<COURSE_CODE>/   (for shared levels, e.g., 100 and 200)
   where <SEM> is "1" or "2". No PQ/GENERAL/Assignments/etc. buckets are ever created.

• **Dept whitelist**: we never create departments beyond an explicit whitelist
  (pass with --dept-whitelist EEE,AAE,PTE,CVL,BME,CHE,COE,....). If omitted, we
  learn the set from courses.json and stick to it.

• **Shared levels**: specify with --shared-levels 100,200 (default). Any course
  whose number starts with one of these is routed to the **shared root** at the
  same level as departments:  <OUTPUT>/<LEVEL>/<SEM>/<COURSE_CODE>/

• **Exact-code first**: if the filename has an exact course code from the
  catalog (e.g., "CVE 309" or "CVE309"), we assign that course immediately.

• **Dept guard**: we never place a file into a department if the course code’s
  department differs. Similarity cannot override this.

• **Two-pass extraction** (up to --pages-max, default 20) with PyMuPDF fallback
  and optional OCR. You can switch to SBERT via --method sbert if installed.

• **Safe fallback**: if we still can’t determine a course, the file is copied to
  <OUTPUT>/_unclassified/ (single folder). Originals are untouched unless --move
  is passed.

Run (example):
  python classifier_v5.py \
    --input-dir "/path/COMPILATION" \
    --courses-json "/path/courses.json" \
    --output-root "/path/SortedDocs" \
    --dept-whitelist EEE,AAE,PTE,CVL,BME,CHE,COE,AAE,BME \
    --shared-levels 100,200 \
    --pages 4 --pages-max 20 \
    --method tfidf --min-score 0.30 --workers 6

MIT License
"""
from __future__ import annotations

import argparse
import concurrent.futures
import dataclasses
import functools
import hashlib
import json
import logging
import os
import re
import shutil
import sys
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# ---------- Optional libs ----------
try:
    from pypdf import PdfReader  # type: ignore
    PYPDF_AVAILABLE = True
except Exception:
    PYPDF_AVAILABLE = False

try:
    import fitz  # PyMuPDF  # type: ignore
    PYMUPDF_AVAILABLE = True
except Exception:
    PYMUPDF_AVAILABLE = False

try:
    import docx  # python-docx  # type: ignore
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation  # python-pptx  # type: ignore
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    import pytesseract  # type: ignore
    PYTESSERACT_AVAILABLE = True
except Exception:
    PYTESSERACT_AVAILABLE = False

try:
    from pdf2image import convert_from_path  # type: ignore
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False

from sklearn.feature_extraction.text import TfidfVectorizer  # type: ignore
from sklearn.metrics.pairwise import cosine_similarity  # type: ignore

SBERT_AVAILABLE = False
try:
    from sentence_transformers import SentenceTransformer  # type: ignore
    import numpy as np  # type: ignore
    SBERT_AVAILABLE = True
except Exception:
    pass

# ---------- Logging ----------
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s", handlers=[logging.StreamHandler(sys.stdout)])
logger = logging.getLogger("doc_classifier")

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt"}
COURSE_CODE_RX = re.compile(r"^[A-Za-z]{2,5}\s?-?\d{2,3}$")
CODE_IN_TEXT_RX = re.compile(r"\b([A-Za-z]{2,5})\s?-?(\d{2,3})\b")

# ---------- Helpers ----------

def normalize_space(s: str) -> str:
    s = s.replace("\x00", " ")
    s = re.sub(r"\s+", " ", s)
    return s.strip()


def safe_segment(name: str) -> str:
    name = re.sub(r"[\\/]+", " ", str(name)).strip()
    name = re.sub(r"[^A-Za-z0-9 _.-]+", "", name)
    name = re.sub(r"\s+", " ", name)
    return name


def semester_to_number(sem: Optional[str]) -> str:
    if not sem: return "1"
    s = str(sem).strip().upper()
    if s in {"1","FIRST","FIRST SEMESTER","SEMESTER 1"}: return "1"
    if s in {"2","SECOND","SECOND SEMESTER","SEMESTER 2"}: return "2"
    return "1"

# ---------- JSON ----------

def read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def write_json(path: Path, obj: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")


def normalize_courses(courses_raw: Any) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []

    def coerce_course(x: Any, hints: Dict[str, Any] | None = None):
        hints = hints or {}
        if isinstance(x, dict):
            c = dict(x)
            for k in ("department_code", "level", "semester"):
                c.setdefault(k, hints.get(k))
            if c.get("code"): c["code"] = str(c["code"]).strip().upper()
            if c.get("title"): c["title"] = str(c["title"]).strip()
            if c.get("department_code"): c["department_code"] = str(c["department_code"]).strip().upper()
            if c.get("level"):
                lev = re.sub(r"\D+", "", str(c["level"]))
                c["level"] = lev or str(c["level"]).strip()
            if c.get("semester"): c["semester"] = str(c["semester"]).strip()
            out.append(c)
        elif isinstance(x, str):
            s = x.strip()
            m = re.match(r"^([A-Za-z]{2,5})\s*-?\s*(\d{2,3})\b(.*)$", s)
            dept = level = sem = None
            code = title = None
            if m:
                dept = m.group(1).upper(); num = m.group(2)
                code = f"{dept} {num}"
                level = num[0] + "00" if len(num) >= 3 else None
                title = m.group(3).strip(" -")
                semu = s.upper()
                if "FIRST" in semu or "1ST" in semu: sem = "FIRST"
                if "SECOND" in semu or "2ND" in semu: sem = "SECOND"
            out.append({
                "code": code or s.upper(),
                "title": title or "",
                "department_code": hints.get("department_code") or dept,
                "level": hints.get("level") or level,
                "semester": hints.get("semester") or sem,
            })

    def walk(node: Any, hints: Dict[str, Any] | None = None):
        hints = hints or {}
        if isinstance(node, list):
            for itm in node: walk(itm, hints)
        elif isinstance(node, dict):
            if all(k in node for k in ("code", "title")):
                coerce_course(node, hints); return
            if all(isinstance(k, str) for k in node.keys()):
                for k, v in node.items():
                    keyu = str(k).upper().strip()
                    if keyu in {"FIRST","SECOND","1","2"}:
                        walk(v, {**hints, "semester": "FIRST" if keyu in {"FIRST","1"} else "SECOND"})
                    elif re.fullmatch(r"\d{3}", keyu):
                        walk(v, {**hints, "level": keyu})
                    elif re.fullmatch(r"[A-Z]{2,5}", keyu):
                        walk(v, {**hints, "department_code": keyu})
                    else:
                        if keyu in {"COURSES","LIST","OUTLINE","CATALOG"}: walk(v, hints)
                        else:
                            if isinstance(v, dict) and ("code" in v or "title" in v):
                                vv = dict(v); vv.setdefault("code", k); coerce_course(vv, hints)
                            else:
                                walk(v, hints)
            else:
                for v in node.values(): walk(v, hints)
        else:
            coerce_course(node, hints)

    walk(courses_raw)

    cleaned: List[Dict[str, Any]] = []
    for c in out:
        code = (c.get("code") or "").strip().upper()
        if not code: continue
        if not COURSE_CODE_RX.match(code):
            continue  # drop headings like "CHEMICAL ENGINEERING"
        if not c.get("department_code"):
            m = re.match(r"^([A-Za-z]{2,5})\s?-?(\d{2,3})", code)
            if m: c["department_code"] = m.group(1).upper()
        cleaned.append(c)
    return cleaned

# ---------- Corpus ----------

def course_to_text(course: Dict[str, Any]) -> str:
    parts: List[str] = []
    for key in ("code","title","department","department_code","level","semester"):
        val = course.get(key)
        if val: parts.append(str(val))
    desc = course.get("description"); topics = course.get("topics") or []
    if desc: parts.append(desc)
    if isinstance(topics, list): parts.extend([str(t) for t in topics])
    elif isinstance(topics, str): parts.append(topics)
    return normalize_space(" ".join(parts))


def build_course_corpus(courses: List[Dict[str, Any]]) -> Tuple[List[str], List[str]]:
    labels: List[str] = []; texts: List[str] = []
    for c in courses:
        code = str(c.get("code") or "").strip()
        if not code: continue
        labels.append(code)
        texts.append(course_to_text(c))
    return labels, texts


def build_code_meta_map(courses: List[Dict[str, Any]]) -> Dict[str, Dict[str, str]]:
    m: Dict[str, Dict[str, str]] = {}
    for c in courses:
        code = str(c.get("code") or "").strip();
        if not code: continue
        dept_code = c.get("department_code") or c.get("dept_code") or ""
        level = c.get("level") or ""; sem = c.get("semester") or ""
        m[code] = {
            "department_code": safe_segment(dept_code) or "UNKNOWN",
            "level": safe_segment(str(level)) or "000",
            "semester_num": semester_to_number(sem),
            "course_code": safe_segment(code),
        }
    return m

# ---------- Extraction ----------

def extract_pdf_pypdf(path: Path, max_pages: int) -> str:
    if not PYPDF_AVAILABLE: return ""
    try:
        reader = PdfReader(str(path))
        pages = min(len(reader.pages), max_pages)
        parts: List[str] = []
        for i in range(pages):
            try:
                t = reader.pages[i].extract_text() or ""
                if t.strip(): parts.append(t)
            except Exception: continue
        return normalize_space(" ".join(parts))
    except Exception:
        return ""


def extract_pdf_pymupdf(path: Path, max_pages: int) -> str:
    if not PYMUPDF_AVAILABLE: return ""
    try:
        parts: List[str] = []
        with fitz.open(str(path)) as doc:
            pages = min(len(doc), max_pages)
            for i in range(pages):
                try:
                    t = doc.load_page(i).get_text("text") or ""
                    if t.strip(): parts.append(t)
                except Exception: continue
        return normalize_space(" ".join(parts))
    except Exception:
        return ""


def extract_text_pdf(path: Path, max_pages: int, ocr: bool=False, ocr_dpi: int=250) -> str:
    text = extract_pdf_pypdf(path, max_pages) or extract_pdf_pymupdf(path, max_pages)
    if ocr and (not text or len(text) < 60) and PDF2IMAGE_AVAILABLE and PYTESSERACT_AVAILABLE:
        try:
            images = convert_from_path(str(path), dpi=ocr_dpi, first_page=1, last_page=max_pages)
            parts = []
            for img in images:
                try:
                    t = pytesseract.image_to_string(img)
                    if t.strip(): parts.append(t)
                except Exception: pass
            if parts: text = (text + "\n" + "\n".join(parts)).strip()
        except Exception:
            pass
    return normalize_space(text)


def extract_text_docx(path: Path, max_pages: int) -> str:
    if not DOCX_AVAILABLE: return ""
    try:
        doc = docx.Document(str(path))
        para_texts = []; chars = 0
        for p in doc.paragraphs:
            t = p.text.strip();
            if not t: continue
            para_texts.append(t); chars += len(t)
            if chars > max_pages * 1800: break
        return normalize_space(" ".join(para_texts))
    except Exception:
        return ""


def extract_text_pptx(path: Path, max_slides: int) -> str:
    if not PPTX_AVAILABLE: return ""
    texts: List[str] = []
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides):
            if i >= max_slides: break
            for shp in slide.shapes:
                try:
                    if hasattr(shp, "text") and shp.text:
                        texts.append(shp.text.strip())
                    if getattr(shp, "has_table", False):
                        for r in shp.table.rows:
                            for c in r.cells: texts.append(c.text.strip())
                except Exception: continue
    except Exception:
        return ""
    return normalize_space(" ".join([t for t in texts if t]))


def extract_text_txt(path: Path, max_chars: int = 12000) -> str:
    try:
        return normalize_space(path.read_text(encoding="utf-8", errors="ignore")[:max_chars])
    except Exception:
        return ""


def extract_text(path: Path, pages: int, ocr: bool) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf": return extract_text_pdf(path, pages, ocr=ocr)
    if ext == ".docx": return extract_text_docx(path, pages)
    if ext == ".pptx": return extract_text_pptx(path, pages)
    if ext == ".txt": return extract_text_txt(path)
    return ""

# ---------- Vectors ----------
@dataclasses.dataclass
class VectorIndex:
    method: str
    labels: List[str]
    model_name: Optional[str] = None
    tfidf: Optional[TfidfVectorizer] = None
    tfidf_matrix: Optional[Any] = None
    sbert_model: Optional[Any] = None
    sbert_vectors: Optional[Any] = None

def build_index_tfidf(course_texts: List[str]) -> Tuple[TfidfVectorizer, Any]:
    vec = TfidfVectorizer(lowercase=True, stop_words="english", ngram_range=(1,2), min_df=1, max_df=0.95)
    X = vec.fit_transform(course_texts)
    return vec, X

def build_index_sbert(course_texts: List[str], model_name: str = "all-MiniLM-L6-v2"):
    if not SBERT_AVAILABLE: raise RuntimeError("sentence-transformers not installed; use --method tfidf")
    model = SentenceTransformer(model_name)
    vectors = model.encode(course_texts, batch_size=32, show_progress_bar=False, normalize_embeddings=True)
    return model, vectors

# ---------- Classification ----------
@dataclasses.dataclass
class Classification:
    file: Path
    course: Optional[str]
    score: float
    reason: str


def filename_exact_code(name: str, valid_codes: set[str]) -> Optional[str]:
    base = Path(name).stem.upper().replace("_"," ")
    # exact token match
    for code in valid_codes:
        if code.upper() in base:
            return code
    # pattern like AAE331
    m = CODE_IN_TEXT_RX.search(base)
    if m:
        candidate = f"{m.group(1).upper()} {m.group(2)}"
        if candidate in valid_codes:
            return candidate
    return None


def choose_by_similarity(text: str, index: VectorIndex) -> List[Tuple[str,float]]:
    if not text.strip():
        return [(index.labels[0], 0.0)] if index.labels else []
    if index.method == "tfidf":
        q = index.tfidf.transform([text])
        sims = cosine_similarity(q, index.tfidf_matrix)[0]
    else:
        qvec = index.sbert_model.encode([text], normalize_embeddings=True)[0]
        sims = (index.sbert_vectors @ qvec)
    ranked = sorted(zip(index.labels, sims), key=lambda x: x[1], reverse=True)
    return [(c, float(s)) for c,s in ranked]


def classify_one(path: Path, pages: int, pages_max: int, ocr: bool, index: VectorIndex, meta_map: Dict[str,Dict[str,str]], min_score: float, dept_whitelist: set[str], shared_levels: set[str]) -> Classification:
    valid_codes = set(index.labels)

    # 0) Filename exact code
    exact = filename_exact_code(path.name, valid_codes)
    if exact:
        return Classification(path, exact, 1.0, "exact_code_in_filename")

    # 1) Two-pass extraction
    t1 = extract_text(path, pages, ocr=ocr)
    ranked = choose_by_similarity(t1, index)
    best_code, best_score = ranked[0]

    if best_score < min_score and pages_max > pages:
        t2 = extract_text(path, pages_max, ocr=ocr)
        ranked2 = choose_by_similarity(t1 + "\n" + t2, index)
        best_code, best_score = ranked2[0]

    # 2) If score sufficient, accept
    if best_score >= min_score:
        # Dept guard
        meta = meta_map.get(best_code, {})
        dept = (meta.get("department_code") or "").upper()
        if dept and dept not in dept_whitelist:
            # disallow accidental creation of new dept folders
            return Classification(path, None, best_score, "dept_not_in_whitelist")
        return Classification(path, best_code, best_score, "vector_match")

    # 3) Unclassified
    return Classification(path, None, best_score, "below_threshold")

# ---------- Routing ----------

def route_target(course: Optional[str], meta_map: Dict[str,Dict[str,str]], output_root: Path, shared_levels: set[str]) -> Path:
    if not course:
        return output_root / "_unclassified"
    meta = meta_map[course]
    dept = meta.get("department_code", "UNKNOWN") or "UNKNOWN"
    level = meta.get("level", "000") or "000"
    sem = meta.get("semester_num", "1") or "1"
    code = meta.get("course_code", course)
    # shared levels go to root/<LEVEL>/<SEM>/<CODE>
    if level in shared_levels:
        return output_root / level / sem / code
    # normal dept layout
    return output_root / dept / level / sem / code

# ---------- Cache ----------

def cache_paths(output_root: Path, method: str) -> Tuple[Path, Path]:
    meta = output_root / ".cache" / f"courses_{method}.meta.json"
    vecs = output_root / ".cache" / f"courses_{method}.npz"
    return meta, vecs


def load_index_from_cache(output_root: Path, method: str, model_name: Optional[str], course_texts: List[str], course_labels: List[str]):
    meta_path, vecs_path = cache_paths(output_root, method)
    if not meta_path.exists(): return None
    try:
        meta = json.loads(meta_path.read_text(encoding="utf-8"))
        if meta.get("method") != method: return None
        if model_name and meta.get("model_name") != model_name: return None
        if meta.get("labels") != course_labels: return None
        idx = VectorIndex(method=method, labels=course_labels, model_name=model_name)
        if method == "tfidf":
            tfidf, X = build_index_tfidf(course_texts)
            idx.tfidf = tfidf; idx.tfidf_matrix = X
            return idx
        else:
            if not SBERT_AVAILABLE: return None
            import numpy as _np
            model, _ = build_index_sbert(course_texts, model_name=model_name or "all-MiniLM-L6-v2")
            arr = _np.load(vecs_path)["vectors"]
            idx.sbert_model = model; idx.sbert_vectors = arr
            return idx
    except Exception:
        return None


def save_index_to_cache(output_root: Path, idx: VectorIndex, course_texts: List[str]) -> None:
    meta_path, vecs_path = cache_paths(output_root, idx.method)
    meta = {"method": idx.method, "labels": idx.labels, "model_name": idx.model_name}
    meta_path.parent.mkdir(parents=True, exist_ok=True)
    meta_path.write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")
    if idx.method == "sbert" and SBERT_AVAILABLE:
        import numpy as _np
        vecs = idx.sbert_vectors or idx.sbert_model.encode(course_texts, batch_size=32, show_progress_bar=False, normalize_embeddings=True)
        _np.savez_compressed(vecs_path, vectors=vecs)

# ---------- CLI ----------

def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Strict course classifier with shared 100/200 roots and dept whitelist.")
    p.add_argument("--input-dir", required=True)
    p.add_argument("--courses-json", required=True)
    p.add_argument("--output-root", required=True)
    p.add_argument("--dept-whitelist", default="", help="Comma-separated dept codes; if empty, derived from courses.json")
    p.add_argument("--shared-levels", default="100,200", help="Comma-separated levels that live at root (shared)")
    p.add_argument("--pages", type=int, default=4)
    p.add_argument("--pages-max", type=int, default=20)
    p.add_argument("--ocr", action="store_true")
    p.add_argument("--method", choices=["tfidf","sbert"], default="tfidf")
    p.add_argument("--sbert-model", default="all-MiniLM-L6-v2")
    p.add_argument("--min-score", type=float, default=0.30)
    p.add_argument("--workers", type=int, default=min(8, (os.cpu_count() or 4)))
    p.add_argument("--move", action="store_true")
    p.add_argument("--dry-run", action="store_true")
    p.add_argument("--include-exts", default=",".join(sorted(SUPPORTED_EXTS)))
    return p.parse_args()

# ---------- Main ----------

def file_sha256(path: Path, chunk: int = 1 << 20) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        while True:
            b = f.read(chunk)
            if not b: break
            h.update(b)
    return h.hexdigest()


def main():
    args = parse_args()
    input_dir = Path(args.input_dir).expanduser().resolve()
    courses_json = Path(args.courses_json).expanduser().resolve()
    output_root = Path(args.output_root).expanduser().resolve()

    include_exts = {e.strip().lower() if e.strip().startswith('.') else '.'+e.strip().lower() for e in args.include_exts.split(',') if e.strip()}

    if not courses_json.exists():
        logger.error("courses.json not found: %s", courses_json); sys.exit(2)

    courses_raw = read_json(courses_json)
    courses = normalize_courses(courses_raw)
    if not courses:
        logger.error("No courses parsed from %s", courses_json); sys.exit(2)

    labels, texts = build_course_corpus(courses)
    meta_map = build_code_meta_map(courses)

    # Dept whitelist
    if args.dept_whitelist.strip():
        dept_whitelist = {d.strip().upper() for d in args.dept_whitelist.split(',') if d.strip()}
    else:
        dept_whitelist = { (c.get('department_code') or '').upper() for c in courses if c.get('department_code') }
    # Always drop empties
    dept_whitelist = {d for d in dept_whitelist if d}

    # Shared levels
    shared_levels = {l.strip() for l in args.shared_levels.split(',') if l.strip()}

    # Build/load vector index
    if args.method == "tfidf":
        idx = VectorIndex(method="tfidf", labels=labels)
        tfidf, X = build_index_tfidf(texts); idx.tfidf = tfidf; idx.tfidf_matrix = X
    else:
        idx = VectorIndex(method="sbert", labels=labels, model_name=args.sbert_model)
        model, vecs = build_index_sbert(texts, model_name=args.sbert_model); idx.sbert_model = model; idx.sbert_vectors = vecs

    # Walk files
    files: List[Path] = []
    for root, _, fnames in os.walk(input_dir):
        for fn in fnames:
            p = Path(root) / fn
            if p.suffix.lower() in include_exts: files.append(p)
    files.sort()
    if not files:
        logger.warning("No input files found under %s", input_dir); return

    if not args.move:
        logger.info("Copy mode: originals remain untouched.")

    results: List[Dict[str, Any]] = []

    fn = functools.partial(
        classify_one,
        pages=args.pages,
        pages_max=args.pages_max,
        ocr=args.ocr,
        index=idx,
        meta_map=meta_map,
        min_score=args.min_score,
        dept_whitelist=dept_whitelist,
        shared_levels=shared_levels,
    )

    start = time.time()
    with concurrent.futures.ThreadPoolExecutor(max_workers=args.workers) as ex:
        for cls in ex.map(fn, files):
            if cls.course:
                dst_dir = route_target(cls.course, meta_map, output_root, shared_levels)
            else:
                dst_dir = output_root / "_unclassified"

            # enforce whitelist for dept layout destinations
            if cls.course:
                meta = meta_map[cls.course]
                lvl = meta.get("level", "000")
                sem = meta.get("semester_num", "1")
                dept = (meta.get("department_code") or "").upper()
                if lvl not in shared_levels and dept not in dept_whitelist:
                    # if somehow a code has an unknown dept, fail safe
                    dst_dir = output_root / "_unclassified"

            logger.info("FILE: %s -> %s  reason=%s  score=%.3f",
                        cls.file.name, dst_dir, cls.reason, cls.score)

            if not args.dry_run:
                try:
                    dst_dir.mkdir(parents=True, exist_ok=True)
                    dst = dst_dir / cls.file.name
                    if dst.exists():
                        prefix = file_sha256(cls.file)[:8]
                        dst = dst_dir / f"{prefix}_{cls.file.name}"
                    if args.move: shutil.move(str(cls.file), str(dst))
                    else: shutil.copy2(str(cls.file), str(dst))
                except Exception as e:
                    logger.error("Failed to place %s: %s", cls.file, e)

            results.append({
                "file": str(cls.file),
                "assigned": cls.course,
                "target": str(dst_dir),
                "score": cls.score,
                "reason": cls.reason,
            })

    elapsed = time.time() - start
    logger.info("Done in %.2fs. %d classified, %d unclassified.",
                elapsed,
                sum(1 for r in results if r["assigned"]),
                sum(1 for r in results if not r["assigned"]))

    write_json(output_root / "classification_report.json", results)

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        logger.warning("Interrupted by user"); sys.exit(130)
